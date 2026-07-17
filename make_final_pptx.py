import collections
import collections.abc
from pptx import Presentation

prs = Presentation('/home/tarun/Documents/IBM/ds-capstone-template-coursera.pptx')
github_url = "https://github.com/tarunkrjaiswal/IBM-Data-Science"

replacements_by_slide = {
    1: {"Enter your name here": "Tarun Kumar Jaiswal"},
    2: {"Add your text here": "The objective of this capstone project is to predict if the SpaceX Falcon 9 first stage will land successfully. We used the SpaceX API and Web Scraping to collect historical launch data. The data was wrangled, and exploratory data analysis was performed using visualization (Matplotlib/Seaborn) and SQL. Interactive visual analytics dashboards were built using Folium and Plotly Dash. Finally, predictive machine learning models were created, identifying the Decision Tree as highly effective, predicting successful landings with great accuracy."},
    3: {"Add your text here": "SpaceX advertises Falcon 9 rocket launches at around $62 million, which is significantly cheaper than competitors. This cost-saving is largely due to their ability to reuse the first stage. Predicting whether the first stage will land successfully is crucial for estimating launch costs. This project seeks to build a machine learning pipeline to determine if the first stage will land based on parameters like payload mass, orbit, launch site, and flight number."},
    8: {"Present your data collection with SpaceX REST calls using key phrases and flowcharts": "Data Collection: Made GET requests to the SpaceX API. Filtered for Falcon 9 launches. Processed the JSON payload to extract booster versions, payloads, launch sites, and outcomes."},
    9: {"Present your web scraping process using key phrases and flowcharts": "Web Scraping: Extracted data from the Wikipedia page for Falcon 9 launches using BeautifulSoup. Parsed HTML tables to extract launch records and constructed a pandas DataFrame."},
    10: {"You need to present your data wrangling process using key phrases and flowcharts": "Data Wrangling: Handled missing values by imputing means (e.g. payload mass). Created binary target labels: 1 for successful landings, 0 for failures. One-hot encoded categorical variables.",
         "Describe how data were processed": "Missing values were handled, categorical data was one-hot encoded, and outcomes were binary labeled."},
    11: {"Summarize what charts were plotted and why you used those charts": "Visualizations: Scatter plots were used to identify relationships between payload, flight number, and launch sites. Bar charts and line graphs were used to evaluate success rates across orbit types and over time."},
    12: {"Using bullet point format, summarize the SQL queries you performed": "SQL Queries performed:\n- Identified unique launch sites\n- Filtered launch records by specific launch sites and boosters\n- Aggregated total and average payload masses\n- Queried records of first successful landings\n- Grouped and ranked landing outcomes"},
    13: {"Summarize what map objects such as markers, circles, lines, etc. you created and added to a folium map": "Folium Maps:\n- Markers: To denote launch site locations on a global map.\n- MarkerClusters: To group launches by site and color code outcomes (Green=Success, Red=Failure).\n- Polylines: To calculate and visualize distances from launch sites to proximities like coastlines and railways.",
         "Explain why you added those objects": "These objects visually demonstrate proximity and success clusters."},
    14: {"Summarize what plots/graphs and interactions you have added to a dashboard": "Plotly Dash Dashboard:\n- Dropdown & Slider: To allow users to filter by launch site and payload range.\n- Pie Chart: Added to visualize the proportion of successful launches per site.\n- Scatter Plot: Added to show the correlation between payload mass and launch outcome.",
         "Explain why you added those plots and interactions": "To provide an interactive tool for exploring success correlations."},
    15: {"Summarize how you built, evaluated, improved, and found the best performing classification model": "Predictive Analysis:\n- Preprocessed data using StandardScaler.\n- Built Classification models: Logistic Regression, SVM, Decision Tree, and KNN.\n- Used GridSearchCV with 10-fold cross-validation to tune hyperparameters.\n- The Decision Tree classifier performed the best on test data.",
         "You need present your model development process using key phrases and flowchart": "Models evaluated: Logistic Regression, SVM, Decision Tree, KNN."},
    18: {"Show the screenshot of the scatter plot with explanations": "Observation: As flight number increases, the success rate also generally improves. Most launches are concentrated at CCAFS and KSC launch sites."},
    19: {"Show the screenshot of the scatter plot with explanations": "Observation: Higher payload masses (e.g. > 10,000 kg) do not negatively impact launch success, and are frequently launched from KSC LC 39A and CCAFS SLC 40."},
    20: {"Show the screenshot of the scatter plot with explanations": "Observation: Orbits like ES-L1, GEO, HEO, and SSO have a 100% success rate, while SO has a 0% success rate."},
    21: {"Show the screenshot of the scatter plot with explanations": "Observation: Higher flight numbers are correlated with successful insertions into various orbits like VLEO."},
    22: {"Show the screenshot of the scatter plot with explanations": "Observation: VLEO orbit is heavily used for payloads around 15,000 kg, predominantly successfully."},
    23: {"Show the screenshot of the scatter plot with explanations": "Observation: The overall launch success rate has been consistently improving over time, climbing sharply from 2013 and maintaining high levels through 2020."},
    
    # SQL queries
    24: {"Present your query result with a short explanation here": "Query Result: CCAFS LC-40, CCAFS SLC-40, KSC LC-39A, VAFB SLC-4E. These are the primary launch facilities for SpaceX Falcon 9."},
    25: {"Present your query result with a short explanation here": "Query Result: Successfully retrieved 5 launch records from sites beginning with 'CCA' (CCAFS LC-40)."},
    26: {"Present your query result with a short explanation here": "Query Result: Total Payload Mass for NASA (CRS) is 45,596 kg."},
    27: {"Present your query result with a short explanation here": "Query Result: The average payload mass for booster version F9 v1.1 is 2,928 kg."},
    28: {"Present your query result with a short explanation here": "Query Result: The first successful ground pad landing was on 2015-12-22."},
    29: {"Present your query result with a short explanation here": "Query Result: Boosters F9 FT B1022, B1026, B1021.2, and B1031.2 successfully landed on drone ships with payload between 4000 and 6000 kg."},
    30: {"Present your query result with a short explanation here": "Query Result: 73 Successful missions, 27 Failed missions (including various modes of failure)."},
    31: {"Present your query result with a short explanation here": "Query Result: Boosters B1048.4, B1049.4, B1051.3, B1056.4 carried the maximum payload of 15,600 kg."},
    32: {"Present your query result with a short explanation here": "Query Result: In 2015, F9 v1.1 B1012 and B1015 failed drone ship landings from CCAFS LC-40."},
    33: {"Present your query result with a short explanation here": "Query Result: Success (drone ship) and Success (ground pad) were the most frequent outcomes between those dates."},

    # Folium
    35: {"Replace <Folium map screenshot 1> title with an appropriate title": "Launch Sites on Global Map",
         "Explain the important elements and findings on the screenshot": "Observation: All launch sites are located near the coastline to ensure safety over the ocean during launches."},
    36: {"Replace <Folium map screenshot 2> title with an appropriate title": "Launch Outcomes on Global Map",
         "Explain the important elements and findings on the screenshot": "Observation: Successful landings (green) and failures (red) show that KSC LC-39A has the highest concentration of successful outcomes."},
    37: {"Replace <Folium map screenshot 3> title with an appropriate title": "Launch Site Proximities",
         "Explain the important elements and findings on the screenshot": "Observation: Launch sites are strategically placed near railways and highways for transportation logistics, and far from densely populated areas."},

    # Dashboard
    39: {"Replace <Dashboard screenshot 1> title with an appropriate title": "Launch Success Pie Chart (All Sites)",
         "Explain the important elements and findings on the screenshot": "Observation: KSC LC-39A contributed to the largest share of successful launches."},
    40: {"Replace <Dashboard screenshot 2> title with an appropriate title": "Launch Success Pie Chart (KSC LC-39A)",
         "Explain the important elements and findings on the screenshot": "Observation: For KSC LC-39A, the vast majority of launches were successful."},
    41: {"Replace <Dashboard screenshot 3> title with an appropriate title": "Payload vs Launch Outcome Scatter Plot",
         "Explain the important elements and findings on the screenshot": "Observation: Payloads between 2000-4000 kg have mixed outcomes, but higher payloads tend to be successfully launched more consistently."},

    # Machine Learning
    43: {"Find which model has the highest classification accuracy": "The Decision Tree model achieved the highest test classification accuracy (83.3%), matching the SVM model."},
    44: {"Show the confusion matrix of the best performing model with an explanation": "The confusion matrix shows the model correctly predicted 12 True Positives (successful landings) and 3 False Positives (predicted success, but actually failed). It effectively distinguishes classes but struggles slightly with false positives."},
    45: {"Point 1\nPoint 2\nPoint 3\nPoint 4\n…": "1. Launch success rates have improved significantly over time, particularly after 2013.\n2. Orbits like ES-L1, GEO, HEO, and SSO have 100% success rates.\n3. KSC LC-39A is the most successful launch site for Falcon 9.\n4. Machine learning classification models can accurately predict landing outcomes, with Decision Tree classifiers providing robust accuracy."}
}

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    
    # Generic replacement
    if slide_num in replacements_by_slide:
        for search, replacement in replacements_by_slide[slide_num].items():
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text
                if search in text:
                    shape.text = text.replace(search, replacement)
    
    # Github URL replacement & cleaning up specific instructions that appear globally
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "Add the GitHub URL" in shape.text:
            shape.text = f"GitHub URL: {github_url}"
        # Some global cleanup
        if "Replace <" in shape.text and "> title with an appropriate title" in shape.text:
            shape.text = "Visual Insights"

prs.save('/home/tarun/Documents/IBM/Data_Science_Capstone_Project_Report_Final.pptx')
print("Saved Data_Science_Capstone_Project_Report_Final.pptx")
