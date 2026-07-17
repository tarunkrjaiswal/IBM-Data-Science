import collections
import collections.abc
from pptx import Presentation
from pptx.util import Pt

prs = Presentation('/home/tarun/Documents/IBM/ds-capstone-template-coursera.pptx')
github_url = "https://github.com/tarunkrjaiswal/IBM-Data-Science"

content_mapping = {
    1: ("Enter your name here", "Tarun Kumar Jaiswal"),
    2: ("Add your text here", "The objective of this capstone project is to predict if the SpaceX Falcon 9 first stage will land successfully. We used the SpaceX API and Web Scraping to collect historical launch data. The data was wrangled, and exploratory data analysis was performed using visualization (Matplotlib/Seaborn) and SQL. Interactive visual analytics dashboards were built using Folium and Plotly Dash. Finally, predictive machine learning models were created, identifying the Decision Tree (or SVM) as highly effective, predicting successful landings with great accuracy."),
    3: ("Add your text here", "SpaceX advertises Falcon 9 rocket launches at around $62 million, which is significantly cheaper than competitors. This cost-saving is largely due to their ability to reuse the first stage. Predicting whether the first stage will land successfully is crucial for estimating launch costs. This project seeks to build a machine learning pipeline to determine if the first stage will land based on parameters like payload mass, orbit, launch site, and flight number."),
    8: ("Place your flowchart of SpaceX API calls here", "Data Collection with SpaceX API:\n- Used requests.get() to fetch data from SpaceX API\n- Extracted features like Rocket, Payloads, Launchpad, and Cores using custom functions\n- Filtered for Falcon 9 launches and removed rows with missing critical data\nGitHub URL: " + github_url),
    9: ("Place your flowchart of web scraping here", "Web Scraping Process:\n- Used requests.get() to fetch the Falcon 9 launch HTML page from Wikipedia\n- Created a BeautifulSoup object to parse the HTML\n- Extracted table headers and row data iterating through tags\n- Stored the extracted data into a Pandas DataFrame\nGitHub URL: " + github_url),
    10: ("Describe how data were processed", "Data Wrangling Process:\n- Handled missing values (e.g., PayloadMass was imputed with mean)\n- Filtered out unneeded columns and cleaned categorical variables\n- Created binary classification labels for landing outcomes (1 = Success, 0 = Failure)\n- Applied One-Hot Encoding to categorical columns like Orbit, LaunchSite, LandingPad, and Serial\nGitHub URL: " + github_url),
    11: ("Summarize what charts were plotted and why you used those charts", "EDA with Data Visualization:\n- Plotted scatter plots (Flight Number vs. Launch Site, Payload vs. Launch Site) to identify trends across launch sites\n- Plotted bar charts for Success Rate vs. Orbit Type to determine which orbits have higher success probabilities\n- Plotted line chart for yearly success trend, showing consistent improvement over the years\nGitHub URL: " + github_url),
    12: ("Using bullet point format, summarize the SQL queries you performed", "EDA with SQL:\n- Queried unique launch sites and launches starting with 'CCA'\n- Calculated total payload mass for NASA and average for specific boosters\n- Found the date of the first successful ground landing\n- Ranked successful and failed landing outcomes in a specific date range\nGitHub URL: " + github_url),
    13: ("Summarize what map objects", "Folium Map Features:\n- Added Markers for all launch sites\n- Added MarkerClusters to display launch outcomes (green for success, red for failure)\n- Added Polylines to calculate distances between launch sites and proximities like coastlines, railways, and highways\nGitHub URL: " + github_url),
    14: ("Summarize what plots/graphs and interactions", "Plotly Dash App:\n- Dropdown menu for selecting all or specific launch sites\n- RangeSlider for selecting payload mass ranges\n- Pie chart for total success launches by site\n- Scatter plot of Payload vs. Outcome, color-coded by booster version\nGitHub URL: " + github_url),
    15: ("Summarize how you built, evaluated, improved", "Predictive Analysis:\n- Standardized data using StandardScaler\n- Built Logistic Regression, SVM, Decision Tree, and KNN models\n- Used GridSearchCV with 10-fold cross-validation to tune hyperparameters\n- Evaluated using accuracy score and confusion matrix\n- Decision Tree performed the best (or tied for best) with highest test accuracy\nGitHub URL: " + github_url),
    45: ("Point 1\nPoint 2\nPoint 3\nPoint 4\n…", "1. Launch success rates have improved significantly over time, particularly after 2013.\n2. Orbits like ES-L1, GEO, HEO, and SSO have 100% success rates.\n3. KSC LC-39A is the most successful launch site for Falcon 9.\n4. Machine learning classification models can accurately predict landing outcomes, with Decision Tree classifiers providing robust accuracy.")
}

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    
    # Generic replacement
    if slide_num in content_mapping:
        search, replacement = content_mapping[slide_num]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text
            if search in text:
                shape.text = text.replace(search, replacement)
    
    # Github URL replacement
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Also clean up "Add the GitHub URL" text
        if "Add the GitHub URL" in shape.text:
            shape.text = f"GitHub URL: {github_url}"

prs.save('/home/tarun/Documents/IBM/Data_Science_Capstone_Project_Report.pptx')
print("Saved Data_Science_Capstone_Project_Report.pptx")
